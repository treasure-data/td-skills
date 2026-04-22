#!/usr/bin/env python3
"""
Treasure AI Slide Generator - Simple Python Example

This example creates a basic Treasure AI branded presentation using python-pptx.

Requirements:
    pip install python-pptx

Usage:
    python simple-python-example.py

Output:
    TreasureAI_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Treasure AI Brand Colors
NAVY = RGBColor(45, 64, 170)
PURPLE = RGBColor(132, 123, 242)
PINK = RGBColor(196, 102, 212)
SKY_BLUE = RGBColor(128, 179, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(102, 102, 102)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)


def add_logo(slide):
    """
    Reserve space for Treasure AI logo (top-left corner).
    Logo area: x=0.2in, y=0.15in, w=2.0in, h=0.35in

    To add actual logo image:
    slide.shapes.add_picture(
        'path/to/treasure-ai-logo.png',
        Inches(0.2), Inches(0.15),
        height=Inches(0.35)
    )
    """
    # Logo space reserved - no text overlay
    # Users can manually add logo image file later
    pass


def add_page_number(slide, page_num, total_pages=None, color=GRAY):
    """
    Add page number to bottom-right corner.

    Args:
        slide: Slide object
        page_num: Current page number
        total_pages: Total page count (optional, shows "3/10" format)
        color: Text color (default: GRAY)
    """
    if total_pages:
        page_text = f"{page_num}/{total_pages}"
    else:
        page_text = str(page_num)

    page_box = slide.shapes.add_textbox(
        Inches(11.5), Inches(6.95), Inches(1.5), Inches(0.4)
    )
    page_frame = page_box.text_frame
    page_frame.text = page_text
    page_para = page_frame.paragraphs[0]
    page_para.font.name = 'Arial'
    page_para.font.size = Pt(12)
    page_para.font.color.rgb = color
    page_para.alignment = PP_ALIGN.RIGHT


def create_title_slide(prs, title, subtitle=""):
    """Create title slide with Navy background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Navy background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    # Reserve logo space (top-left)
    add_logo(slide)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.94), Inches(2.37), Inches(11.45), Inches(1.64)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = 'Arial'
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Subtitle (optional)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.94), Inches(4.24), Inches(11.45), Inches(0.75)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.name = 'Arial'
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = WHITE


def create_section_divider(prs, section_title):
    """Create section divider slide with Purple background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Purple background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PURPLE

    # Reserve logo space (top-left)
    add_logo(slide)

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(0.94), Inches(2.81), Inches(11.45), Inches(1.64)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = 'Arial'
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE


def create_content_slide(prs, title, bullet_points):
    """Create content slide with white background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Logo
    add_logo(slide)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.94), Inches(0.75), Inches(11.45), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = 'Arial'
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY

    # Content area with bullet points
    content_box = slide.shapes.add_textbox(
        Inches(0.94), Inches(1.88), Inches(11.45), Inches(4.69)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i > 0:
            content_frame.add_paragraph()
        p = content_frame.paragraphs[i]
        p.text = point
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.level = 0


def create_two_column_slide(prs, title, left_content, right_content):
    """Create 2-column layout slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Logo
    add_logo(slide)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.94), Inches(0.75), Inches(11.45), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = 'Arial'
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.94), Inches(1.88), Inches(5.47), Inches(4.69)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.text = left_content
    left_para = left_frame.paragraphs[0]
    left_para.font.name = 'Arial'
    left_para.font.size = Pt(16)
    left_para.font.color.rgb = BLACK

    # Right column (visual placeholder)
    right_box = slide.shapes.add_textbox(
        Inches(6.92), Inches(1.88), Inches(5.47), Inches(4.69)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.text = right_content
    right_para = right_frame.paragraphs[0]
    right_para.font.name = 'Arial'
    right_para.font.size = Pt(16)
    right_para.font.color.rgb = BLACK


def create_end_slide(prs):
    """Create Thank You slide with Navy background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Navy background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    # Reserve logo space (top-left)
    add_logo(slide)

    # Thank You text
    thank_you_box = slide.shapes.add_textbox(
        Inches(0.94), Inches(2.81), Inches(11.45), Inches(1.64)
    )
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = 'Thank You'
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.name = 'Arial'
    thank_you_para.font.size = Pt(44)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = WHITE
    thank_you_para.alignment = PP_ALIGN.CENTER


# Build presentation
print("Creating Treasure AI presentation...")

# Slide 1: Title
create_title_slide(
    prs,
    title="Treasure AI Platform Overview",
    subtitle="Empowering Data-Driven Marketing"
)

# Slide 2: Section Divider
create_section_divider(prs, "What is Treasure AI?")

# Slide 3: Content
create_content_slide(
    prs,
    title="Customer Data Platform",
    bullet_points=[
        "Unified customer data across all touchpoints",
        "Real-time personalization and segmentation",
        "Privacy-first data management",
        "Enterprise-grade security and compliance"
    ]
)

# Slide 4: 2-Column Layout
create_two_column_slide(
    prs,
    title="Key Benefits",
    left_content="• Increase customer engagement\n• Improve conversion rates\n• Reduce operational costs\n• Accelerate time to market",
    right_content="[Visual: Chart showing ROI growth over 12 months]"
)

# Slide 5: Section Divider
create_section_divider(prs, "Platform Features")

# Slide 6: Content
create_content_slide(
    prs,
    title="Core Capabilities",
    bullet_points=[
        "Data Collection: Web, mobile, server-side tracking",
        "Data Integration: 200+ pre-built connectors",
        "Audience Builder: Drag-and-drop segmentation",
        "Activation: Multi-channel campaign orchestration"
    ]
)

# Slide 7: End
create_end_slide(prs)

# Optional: Add page numbers to all slides (except title and end slides)
# Uncomment the following lines to enable page numbering:
#
# total_slides = len(prs.slides)
# for idx, slide in enumerate(prs.slides):
#     # Skip title slide (0) and end slide (last)
#     if idx > 0 and idx < total_slides - 1:
#         add_page_number(slide, idx, total_slides - 2, color=GRAY)

# Save presentation
output_file = 'TreasureAI_Presentation.pptx'
prs.save(output_file)
print(f"✓ Presentation created: {output_file}")
print(f"  Total slides: {len(prs.slides)}")
print(f"  Size: 16:9 (13.33\" x 7.5\")")
print(f"\nTo add page numbers, uncomment lines 307-312 and re-run.")
